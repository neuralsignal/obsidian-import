"""PPTX text extraction using python-pptx.

Extracts per-slide text, speaker notes, tables, and embedded images.
"""

from __future__ import annotations

import logging
from collections.abc import Callable
from pathlib import Path
from typing import TYPE_CHECKING

from obsidian_import.config import MediaConfig

if TYPE_CHECKING:
    from pptx.table import Table

from obsidian_import.exceptions import ExtractionError
from obsidian_import.extraction_result import ExtractionResult, MediaFile
from obsidian_import.formatting import make_media_wikilink, render_markdown_table
from obsidian_import.media import attempt_save_image, generate_media_filename
from obsidian_import.timeout import run_with_timeout

log = logging.getLogger(__name__)

_PICTURE_SHAPE_TYPE = 13


def extract(path: Path, timeout_seconds: int, isolation: str, media_config: MediaConfig) -> ExtractionResult:
    """Extract text and images from a PPTX file, returning ExtractionResult."""
    return run_with_timeout(_extract_pptx, (path, media_config), timeout_seconds, "PPTX", path, isolation)


def _extract_pptx(path: Path, media_config: MediaConfig) -> ExtractionResult:
    """Internal PPTX extraction logic."""
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation(str(path))
    sections: list[str] = [f"# {path.stem}"]
    media_files: list[MediaFile] = []

    slide_width = prs.slide_width
    slide_height = prs.slide_height
    if slide_width and slide_height:
        w_in = slide_width / Inches(1)
        h_in = slide_height / Inches(1)
        sections.append(f'*Slide dimensions: {w_in:.1f}" x {h_in:.1f}"*')

    for i, slide in enumerate(prs.slides, 1):
        slide_sections: list[str] = []

        title = ""
        if slide.shapes.title:
            title = slide.shapes.title.text.strip()

        if title:
            slide_sections.append(f"## Slide {i}: {title}")
        else:
            slide_sections.append(f"## Slide {i}")

        body_texts, slide_media = _extract_slide_content(
            slide,
            i,
            path,
            media_config,
        )
        media_files.extend(slide_media)

        if body_texts:
            slide_sections.append("\n".join(body_texts))

        if slide.has_notes_slide:
            notes_frame = slide.notes_slide.notes_text_frame
            notes_text = notes_frame.text.strip() if notes_frame else ""
            if notes_text:
                slide_sections.append(f"\n> **Speaker Notes:** {notes_text}")

        if len(slide_sections) > 1:
            sections.append("\n\n".join(slide_sections))
        else:
            sections.append(slide_sections[0])

    return ExtractionResult(
        markdown="\n\n".join(sections),
        media_files=tuple(media_files),
    )


def _extract_slide_content(
    slide: object,
    slide_number: int,
    path: Path,
    media_config: MediaConfig,
) -> tuple[list[str], list[MediaFile]]:
    """Extract text, tables, and images from all shapes on a single slide."""
    title_shape = slide.shapes.title  # type: ignore[attr-defined]
    body_texts = _extract_text_from_shapes(slide, title_shape)
    body_texts.extend(_extract_tables_from_shapes(slide))
    image_texts, media_files = _extract_images_from_shapes(slide, slide_number, path, media_config)
    body_texts.extend(image_texts)
    return body_texts, media_files


def _extract_text_from_shapes(slide: object, title_shape: object | None) -> list[str]:
    """Extract text content from all text-frame shapes, skipping the title shape."""
    texts: list[str] = []
    for shape in slide.shapes:  # type: ignore[attr-defined]
        if not shape.has_text_frame or shape == title_shape:
            continue
        for para in shape.text_frame.paragraphs:
            text = para.text.strip()
            if text:
                level = para.level or 0
                if level > 0:
                    texts.append(f"{'  ' * level}- {text}")
                else:
                    texts.append(text)
    return texts


def _extract_tables_from_shapes(slide: object) -> list[str]:
    """Extract markdown tables from all table shapes."""
    tables: list[str] = []
    for shape in slide.shapes:  # type: ignore[attr-defined]
        if shape.has_table:
            table_md = _extract_table(shape.table)
            if table_md:
                tables.append(table_md)
    return tables


def _extract_images_from_shapes(
    slide: object,
    slide_number: int,
    path: Path,
    media_config: MediaConfig,
) -> tuple[list[str], list[MediaFile]]:
    """Extract images from picture shapes, returning wikilinks and media files."""
    texts: list[str] = []
    media_files: list[MediaFile] = []
    image_index = 0
    for shape in slide.shapes:  # type: ignore[attr-defined]
        if not (media_config.extract_images and shape.shape_type == _PICTURE_SHAPE_TYPE):
            continue
        image_index += 1
        try:
            ext = _mime_to_extension(shape.image.content_type)
        except (AttributeError, ValueError):
            ext = ".png"
        filename = generate_media_filename(f"slide{slide_number}", image_index, ext)
        mf = attempt_save_image(
            _make_pptx_image_reader(shape, slide_number),
            filename,
            media_config,
            f"slide {slide_number} of {path}",
        )
        if mf is not None:
            media_files.append(mf)
            texts.append(make_media_wikilink(path.stem, mf.filename))
    return texts, media_files


def _make_pptx_image_reader(shape: object, slide_number: int) -> Callable[[], bytes]:
    """Return a callable that reads image bytes from a PPTX shape."""

    def _read() -> bytes:
        try:
            return shape.image.blob  # type: ignore[union-attr]
        except (AttributeError, ValueError) as exc:
            raise ExtractionError(f"PPTX image on slide {slide_number} unavailable: {exc}") from exc

    return _read


def _mime_to_extension(content_type: str) -> str:
    """Convert MIME type to file extension."""
    mime_map = {
        "image/png": ".png",
        "image/jpeg": ".jpeg",
        "image/gif": ".gif",
        "image/bmp": ".bmp",
        "image/tiff": ".tiff",
        "image/svg+xml": ".svg",
        "image/webp": ".webp",
    }
    return mime_map.get(content_type, ".png")


def _extract_table(table: Table) -> str:
    """Extract a PPTX table as markdown."""
    rows: list[list[str]] = []
    for row in table.rows:
        cells = [cell.text.strip().replace("\n", " ") for cell in row.cells]
        rows.append(cells)

    if not rows:
        return ""

    return render_markdown_table(rows)
